#!/usr/bin/env python3
import argparse, hashlib, html, json, os, re, shutil, subprocess, tempfile
from collections import defaultdict
from pathlib import Path
from urllib.parse import quote
import requests

DOC_EXTS = {'.pdf','.docx','.pptx','.odt','.rtf','.tex','.md','.txt','.html','.htm'}
ADS_RE = re.compile(r'(?<![A-Za-z])ADS(?![A-Za-z])', re.I)
DS_RE = re.compile(r'\b(?:applied\s+data\s+science|data\s+science)\b', re.I)
AUTHOR_RE = re.compile(r'\b(?:Brian\s+(?:[A-Z]\.?\s+)?Tenneson|B\.?\s*Tenneson|Tenneson,\s*Brian)\b', re.I)

CATEGORIES = [
    ('healthcare_biomedical', 'Healthcare & Biomedical Analytics', [
        'cancer','tumor','diabetes','heart disease','cardiovascular','lung','kidney','renal','clinical','patient','medical','health','disease'
    ]),
    ('sports_performance', 'Sports & Performance Analytics', [
        'nba','basketball','sports','player','team','season','game','athlete'
    ]),
    ('nlp_text', 'NLP & Text Analytics', [
        'natural language','nlp','spam','email','text classification','token','sentiment','document classification'
    ]),
    ('statistical_methods', 'Statistical Methods & Experimental Design', [
        'correlation','regression','hypothesis','anova','probability','confidence interval','sampling','experiment','statistical','p-value','chi-square'
    ]),
    ('machine_learning', 'Machine Learning & Predictive Modeling', [
        'machine learning','random forest','decision tree','classifier','classification','logistic regression','neural network','prediction','predictive','supervised learning','cross-validation'
    ]),
    ('data_preparation_eda', 'Data Preparation, Cleaning & EDA', [
        'data cleaning','cleaning','preprocessing','missing data','missing value','data wrangling','feature engineering','exploratory data analysis','eda','outlier'
    ]),
]

METHOD_TAGS = {
    'Regression': ['regression','logistic regression','linear regression'],
    'Classification': ['classification','classifier'],
    'Random Forest': ['random forest'],
    'Visualization': ['visualization','visualisation','plot','chart','dashboard'],
    'Correlation': ['correlation'],
    'Data Cleaning': ['data cleaning','cleaning','preprocessing'],
    'Experimental Design': ['experimental design','experiment'],
    'Healthcare': ['cancer','tumor','diabetes','heart disease','lung','kidney','clinical','patient','health'],
    'Sports': ['nba','basketball','sports','athlete'],
    'NLP': ['natural language','nlp','spam','text classification'],
}

def api_headers(token):
    h={'Accept':'application/vnd.github+json','X-GitHub-Api-Version':'2022-11-28','User-Agent':'ads-wing-builder'}
    if token: h['Authorization']=f'Bearer {token}'
    return h

def gh_get(session, url, token, **kwargs):
    r=session.get(url, headers=api_headers(token), timeout=60, **kwargs)
    r.raise_for_status(); return r

def list_public_repos(session, owner, token):
    repos=[]; page=1
    while True:
        data=gh_get(session, f'https://api.github.com/users/{owner}/repos', token,
                    params={'per_page':100,'page':page,'type':'owner','sort':'full_name'}).json()
        if not data: break
        repos.extend([r for r in data if not r.get('private') and not r.get('fork')])
        if len(data)<100: break
        page += 1
    return repos

def list_tree(session, owner, repo, branch, token):
    url=f'https://api.github.com/repos/{owner}/{repo}/git/trees/{quote(branch, safe="")}'
    data=gh_get(session,url,token,params={'recursive':'1'}).json()
    return data.get('tree',[])

def download(session, owner, repo, branch, path, token):
    url=f'https://raw.githubusercontent.com/{owner}/{repo}/{quote(branch, safe="")}/{quote(path, safe="/")}'
    return gh_get(session,url,token).content

def decode_text(data):
    for enc in ('utf-8','utf-8-sig','cp1252','latin1'):
        try: return data.decode(enc)
        except UnicodeDecodeError: pass
    return data.decode('utf-8',errors='replace')

def extract_text(path: Path):
    ext=path.suffix.lower()
    if ext in {'.tex','.md','.txt','.html','.htm'}:
        return decode_text(path.read_bytes())
    if ext=='.pdf':
        out=path.with_suffix('.txt')
        subprocess.run(['pdftotext','-layout',str(path),str(out)],check=True,stdout=subprocess.DEVNULL,stderr=subprocess.DEVNULL)
        return out.read_text(errors='replace') if out.exists() else ''
    if ext=='.docx':
        from docx import Document
        d=Document(str(path)); return '\n'.join(p.text for p in d.paragraphs)
    if ext=='.pptx':
        from pptx import Presentation
        p=Presentation(str(path)); chunks=[]
        for s in p.slides:
            for sh in s.shapes:
                if hasattr(sh,'text'): chunks.append(sh.text)
        return '\n'.join(chunks)
    if ext=='.odt':
        from zipfile import ZipFile
        from bs4 import BeautifulSoup
        with ZipFile(path) as z: raw=z.read('content.xml')
        return BeautifulSoup(raw,'xml').get_text('\n')
    if ext=='.rtf':
        from striprtf.striprtf import rtf_to_text
        return rtf_to_text(decode_text(path.read_bytes()))
    return ''

def normalize(s): return re.sub(r'\s+',' ',s).strip()

def title_from_text(text, filename):
    lines=[normalize(x) for x in text.splitlines() if normalize(x)]
    skip=re.compile(r'^(?:Brian\s+Tenneson|Tenneson,\s*Brian|ADS\s*\d*|Applied Data Science)$',re.I)
    for line in lines[:40]:
        if 6 <= len(line) <= 150 and not skip.match(line):
            return line
    return Path(filename).stem.replace('_',' ').strip()

def category_for(text):
    t=text.lower()
    scored=[]
    for slug,label,terms in CATEGORIES:
        distinct=sum(1 for term in terms if term.lower() in t)
        occurrences=sum(t.count(term.lower()) for term in terms)
        scored.append((distinct,occurrences,slug,label))
    score_map={slug:{'distinct_terms':d,'occurrences':o} for d,o,slug,_ in scored if d}
    for preferred in ('healthcare_biomedical','sports_performance','nlp_text'):
        hit=next((x for x in scored if x[2]==preferred),None)
        if hit and hit[0] >= 1:
            return hit[2], hit[3], score_map
    scored.sort(key=lambda x:(x[0],x[1]), reverse=True)
    if scored and scored[0][0] > 0:
        return scored[0][2], scored[0][3], score_map
    return 'general', 'General Applied Data Science', {}

def tags_for(text):
    t=text.lower(); tags=[]
    for tag,terms in METHOD_TAGS.items():
        if any(term.lower() in t for term in terms): tags.append(tag)
    return tags

def safe_name(name):
    s=re.sub(r'[^A-Za-z0-9._-]+','_',name).strip('_')
    return s or 'document'

def evidence_snippet(text, match, radius=90):
    if not match: return ''
    a=max(0,match.start()-radius); b=min(len(text),match.end()+radius)
    return normalize(text[a:b])[:240]

def author_evidence(text,path):
    m=AUTHOR_RE.search(text)
    if m: return True, evidence_snippet(text,m)
    if re.search(r'tenneson',path,re.I): return True, 'Filename/path contains “Tenneson”.'
    return False,''

def write_manifest_md(records, out):
    lines=['# Applied Data Science (ADS) Wing Manifest','',
           'Generated from public repositories owned by `btenneson`. Eligibility requires author evidence plus a standalone `ADS` marker and `data science` or `applied data science` in extracted document text.','']
    grouped=defaultdict(list)
    for r in records: grouped[r['category']].append(r)
    for label in sorted(grouped):
        lines += [f'## {label}','']
        for r in sorted(grouped[label],key=lambda x:x['title'].lower()):
            lines.append(f"- **{r['title']}** — `{r['local_path']}`")
            lines.append(f"  - Source: {r['sources'][0]['repo']}/{r['sources'][0]['path']}")
            lines.append(f"  - Tags: {', '.join(r['tags']) if r['tags'] else '—'}")
        lines.append('')
    out.write_text('\n'.join(lines),encoding='utf-8')

def build_html(records, generated_at):
    grouped=defaultdict(list)
    for r in records: grouped[r['category']].append(r)
    cards=[]
    for category in sorted(grouped):
        items=[]
        for r in sorted(grouped[category],key=lambda x:x['title'].lower()):
            rel=quote(r['web_path'],safe='/')
            src=html.escape(r['sources'][0]['html_url'],quote=True)
            tags=', '.join(r['tags']) if r['tags'] else 'Applied Data Science'
            items.append(f'<article class="card"><h3>{html.escape(r["title"])}</h3><p>{html.escape(tags)}</p><p class="links"><a href="{rel}">Open document</a><a href="{src}">Original GitHub source</a></p></article>')
        cards.append(f'<section><h2>{html.escape(category)}</h2><div class="grid">{"".join(items)}</div></section>')
    return f'''<!doctype html><html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1"><title>Applied Data Science (ADS) — Brian Tenneson</title><style>body{{font-family:system-ui;line-height:1.45;margin:0}}main{{max-width:1100px;margin:auto;padding:24px}}.grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(280px,1fr));gap:12px}}.card{{border:1px solid #8886;border-radius:10px;padding:12px}}.card h3{{font-size:1rem;margin:0 0 8px}}.links{{display:flex;gap:12px;flex-wrap:wrap}}input{{font:inherit;padding:9px;width:min(100%,460px);margin:8px 0 18px}}section{{margin:28px 0}}code{{overflow-wrap:anywhere}}</style></head><body><main><p><a href="../">← Main library</a></p><h1>Applied Data Science (ADS)</h1><p>Works by Brian Tenneson collected from public GitHub repositories when the extracted document text identifies ADS / Applied Data Science context. Documents are grouped by their primary subject matter; originals remain in place.</p><p><a href="MANIFEST.md">Manifest</a> · <a href="manifest.json">Machine-readable manifest</a></p><input id="q" type="search" placeholder="Search ADS wing…"><div id="g">{''.join(cards)}</div><p><small>Manifest generated {html.escape(generated_at)}.</small></p></main><script>q.oninput=()=>document.querySelectorAll('.card').forEach(x=>x.hidden=!x.innerText.toLowerCase().includes(q.value.toLowerCase()))</script></body></html>'''

def patch_main_index(path: Path):
    s=path.read_text(encoding='utf-8')
    if 'href="ADS/"' in s: return
    s=s.replace('Stable browser-reader links for all 25 papers represented in this repository.',
                'Stable browser-reader links for the core paper library, plus specialized subject wings.')
    needle='<input id="q" type="search" placeholder="Search papers…">'
    wing='<div class="grid" style="margin:12px 0 24px"><article class="card"><h3>Applied Data Science (ADS)</h3><p>Graduate-school and related applied data-science work, grouped by subject matter.</p><p class="links"><a href="ADS/">Browse ADS wing</a><a href="ADS/MANIFEST.md">Manifest</a></p></article></div>'
    if needle in s: s=s.replace(needle,wing+needle,1)
    path.write_text(s,encoding='utf-8')

def main():
    ap=argparse.ArgumentParser()
    ap.add_argument('--owner',default='btenneson')
    ap.add_argument('--root',default='.')
    ap.add_argument('--token',default=os.getenv('GITHUB_TOKEN',''))
    ap.add_argument('--generated-at',default=os.getenv('GITHUB_RUN_STARTED_AT',''))
    args=ap.parse_args()
    root=Path(args.root).resolve(); archive=root/'ADS'; web=root/'docs'/'ADS'
    for d in (archive,web):
        if d.exists(): shutil.rmtree(d)
        d.mkdir(parents=True)
    session=requests.Session(); records=[]; by_hash={}
    repos=list_public_repos(session,args.owner,args.token)
    for repo in repos:
        name=repo['name']
        if name.lower()=='pub': continue
        branch=repo['default_branch']
        try: tree=list_tree(session,args.owner,name,branch,args.token)
        except Exception as e:
            print(f'WARN tree {name}: {e}'); continue
        for item in tree:
            if item.get('type')!='blob': continue
            path=item['path']; ext=Path(path).suffix.lower()
            if ext not in DOC_EXTS: continue
            try: data=download(session,args.owner,name,branch,path,args.token)
            except Exception as e:
                print(f'WARN download {name}/{path}: {e}'); continue
            with tempfile.TemporaryDirectory() as td:
                f=Path(td)/('source'+ext); f.write_bytes(data)
                try: text=extract_text(f)
                except Exception as e:
                    print(f'WARN extract {name}/{path}: {e}'); continue
            am=ADS_RE.search(text); dm=DS_RE.search(text); authored,aev=author_evidence(text,path)
            if not (authored and am and dm): continue
            digest=hashlib.sha256(data).hexdigest()
            src={'repo':name,'path':path,'branch':branch,'blob_sha':item.get('sha',''),
                 'html_url':f'https://github.com/{args.owner}/{name}/blob/{quote(branch,safe="")}/{quote(path,safe="/")}',
                 'raw_url':f'https://raw.githubusercontent.com/{args.owner}/{name}/{quote(branch,safe="")}/{quote(path,safe="/")}'}
            if digest in by_hash:
                by_hash[digest]['sources'].append(src); continue
            slug,label,scores=category_for(text)
            filename=safe_name(Path(path).name)
            dest=archive/slug/filename; wdest=web/slug/filename
            if dest.exists():
                filename=safe_name(f'{name}_{Path(path).name}'); dest=archive/slug/filename; wdest=web/slug/filename
            dest.parent.mkdir(parents=True,exist_ok=True); wdest.parent.mkdir(parents=True,exist_ok=True)
            dest.write_bytes(data); wdest.write_bytes(data)
            rec={'title':title_from_text(text,Path(path).name),'category':label,'category_slug':slug,
                 'category_scores':scores,'tags':tags_for(text),'sha256':digest,'bytes':len(data),
                 'local_path':dest.relative_to(root).as_posix(),'web_path':wdest.relative_to(web).as_posix(),
                 'sources':[src],
                 'match_evidence':{'author':aev,'ADS':evidence_snippet(text,am),'data_science':evidence_snippet(text,dm)}}
            records.append(rec); by_hash[digest]=rec
            print(f'MATCH {name}/{path} -> {rec["local_path"]}')
    records.sort(key=lambda r:(r['category'],r['title'].lower()))
    generated=args.generated_at or 'GitHub Actions run'
    payload={'schema_version':1,'owner':args.owner,'generated_at':generated,'document_count':len(records),'documents':records}
    for d in (archive,web):
        (d/'manifest.json').write_text(json.dumps(payload,indent=2,ensure_ascii=False),encoding='utf-8')
        write_manifest_md(records,d/'MANIFEST.md')
    (archive/'README.md').write_text('# Applied Data Science (ADS)\n\nCanonical copied archive for the ADS wing. See `MANIFEST.md` for provenance and categorization.\n',encoding='utf-8')
    (web/'index.html').write_text(build_html(records,generated),encoding='utf-8')
    patch_main_index(root/'docs'/'index.html')
    print(f'Generated {len(records)} unique ADS documents from {len(repos)} public repositories.')

if __name__=='__main__': main()
